"""
utils/ppt_generator.py  

"""

import io
import os
import re
from copy import deepcopy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree


TEMPLATE_PATH = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "inputs", "sample_ppt_template.pptx"
)

BLACK = RGBColor(0x00, 0x00, 0x00)
RED   = RGBColor(0xFF, 0x00, 0x00)


# Helpers

def _is_empty(val) -> bool:
    if val is None:
        return True
    return str(val).strip() in (
        "", "null", "None", "N/A", "n/a", "nan",
        "Information Missing", "none", "Not specified",
    )


def _get(data: dict, *keys, default: str = "") -> str:
    """Return first non-empty value found among the given keys."""
    for k in keys:
        v = data.get(k)
        if not _is_empty(v):
            return str(v).strip()
    return default


# Data mapper — converts any input shape into flat PPT keys

def _build_template_dict(candidate_data: dict) -> dict:
    """
    Accepts candidate_data in ANY of the 4 formats (or a mix) and returns
    a flat dict with the keys _populate_slide1 / _populate_project_slide need.
    """
    d = candidate_data or {}

    # Slide 1 fields 
    name    = _get(d, "name", "NAME", "FULL_NAME")
    role    = _get(d, "current_role", "ROLE", "CURRENT_ROLE")
    summary = _get(d, "objective", "PROFESSIONAL_SUMMARY", "summary")

    # Skills — list or comma string
    ts = d.get("tech_stack") or d.get("TECHNICAL_SKILLS") or ""
    if isinstance(ts, list):
        skills = ", ".join(str(s) for s in ts if s)
    else:
        skills = str(ts).strip()

    # Education — try every key variant
    edu = _get(d, "education", "EDUCATION_DETAILS", "HIGHEST_EDUCATION")
    if _is_empty(edu):
        parts = [
            _get(d, "HIGHEST_EDUCATION"),
            _get(d, "COLLEGE_NAME"),
            _get(d, "EDUCATION_DATES"),
        ]
        edu = " | ".join(p for p in parts if p)

    out = {
        "FULL_NAME":            name,
        "CURRENT_ROLE":         role,
        "PROFESSIONAL_SUMMARY": summary,
        "TECHNICAL_SKILLS":     skills,
        "EDUCATION_DETAILS":    edu,
    }

    # Project fields 
    # Check for Format D: key_projects list
    projects = d.get("key_projects") or []
    using_key_projects = (
        isinstance(projects, list)
        and any(
            isinstance(p, dict) and not _is_empty(p.get("title"))
            for p in projects
        )
    )

    for i in range(1, 5):

        if using_key_projects:
            # Format D: key_projects list 
            proj = (
                projects[i - 1]
                if i <= len(projects) and isinstance(projects[i - 1], dict)
                else {}
            )
            title    = str(proj.get("title",       "") or "").strip()
            duration = str(proj.get("duration",    "") or "").strip()
            prole    = str(proj.get("role",        "") or "").strip()
            desc     = str(proj.get("description", "") or "").strip()
            resps    = proj.get("responsibilities", [])
            if isinstance(resps, list):
                resp_lines = [str(r).strip() for r in resps if str(r).strip()]
            else:
                resp_lines = [str(resps).strip()] if resps else []

        else:
            # Formats A / B / C: flat keys 

            # Project name — same across all formats
            title = _get(d, f"PROJECT{i}_NAME", default="")

            # Duration — Format A: DURATION_PROJECT{i}, Format C: PROJECT{i}_DURATION
            duration = _get(
                d,
                f"PROJECT{i}_DURATION",   # Format C (template_mapper output)
                f"DURATION_PROJECT{i}",   # Format A (test_ppt_generator / legacy)
                default="",
            )

            # Role — Format A: ROLE_PROJECT{i}, Format C: PROJECT{i}_ROLE
            prole = _get(
                d,
                f"PROJECT{i}_ROLE",       # Format C
                f"ROLE_PROJECT{i}",       # Format A
                default="",
            )

            # Description — Format A: Project{i}_Description, Format C: PROJECT{i}_DESCRIPTION
            desc = _get(
                d,
                f"PROJECT{i}_DESCRIPTION",   # Format C
                f"Project{i}_Description",   # Format A  (mixed case!)
                f"project{i}_description",   # lowercase variant
                default="",
            )

            # Responsibilities — try in priority order:
            #  1. Format C: PROJECT{i}_RESPONSIBILITIES (pre-joined string)
            #  2. Format A: project{i}_bullets (list)
            #  3. Format B: ABOUT_PROJECT_BULLET_1..6 (project 1) or PROJECT{i}_BULLET_1..5
            resp_lines = []

            resp_str = _get(d, f"PROJECT{i}_RESPONSIBILITIES", default="")
            if resp_str:
                # Split on newlines, strip bullet chars
                resp_lines = [
                    re.sub(r"^[•\-\*]\s*", "", line).strip()
                    for line in resp_str.split("\n")
                    if line.strip()
                ]
            else:
                # Format A: project{i}_bullets list
                bullets_list = d.get(f"project{i}_bullets") or []
                if isinstance(bullets_list, list) and bullets_list:
                    resp_lines = [str(b).strip() for b in bullets_list if str(b).strip()]
                else:
                    # Format B: individual bullet keys
                    if i == 1:
                        prefix, rng = "ABOUT_PROJECT_BULLET_", range(1, 7)
                    else:
                        prefix, rng = f"PROJECT{i}_BULLET_", range(1, 6)
                    resp_lines = [
                        _get(d, f"{prefix}{b}")
                        for b in rng
                        if not _is_empty(_get(d, f"{prefix}{b}"))
                    ]

        out[f"PROJECT{i}_NAME"]             = title
        out[f"PROJECT{i}_DURATION"]         = duration
        out[f"PROJECT{i}_ROLE"]             = prole
        out[f"PROJECT{i}_DESCRIPTION"]      = desc
        out[f"PROJECT{i}_RESPONSIBILITIES"] = resp_lines   # always a list

    return out

# Shape-level helpers

def _set_para_text(para, text: str, color: RGBColor = BLACK):
    """Overwrite paragraph text; preserve all other formatting."""
    if not para.runs:
        run = para.add_run()
        run.text = text
        try:
            run.font.color.rgb = color
            run.font.bold = False
        except Exception:
            pass
        return

    para.runs[0].text = text
    try:
        para.runs[0].font.color.rgb = color
        para.runs[0].font.bold = False
    except Exception:
        pass
    for run in para.runs[1:]:
        run.text = ""


def _clear_para(para):
    for run in para.runs:
        run.text = ""


def _find_resp_label_idx(tf) -> int:
    for i, para in enumerate(tf.paragraphs):
        if "Responsibilities" in para.text:
            return i
    return -1


def _append_para(tf, text: str):
    """Append a new paragraph by cloning the last one (for overflow bullets)."""
    last_p = tf.paragraphs[-1]._p
    new_p  = deepcopy(last_p)
    for r in new_p.findall(qn("a:r")):
        new_p.remove(r)
    r_el = etree.SubElement(new_p, qn("a:r"))
    rPr  = etree.SubElement(r_el,  qn("a:rPr"))
    rPr.set("lang", "en-US")
    rPr.set("dirty", "0")
    t_el = etree.SubElement(r_el,  qn("a:t"))
    t_el.text = text
    tf._txBody.append(new_p)

# Slide populators

def _populate_slide1(slide, d: dict):
    shapes = slide.shapes

    name    = d.get("FULL_NAME",            "")
    role    = d.get("CURRENT_ROLE",         "")
    summary = d.get("PROFESSIONAL_SUMMARY", "")
    edu     = d.get("EDUCATION_DETAILS",    "")
    skills  = d.get("TECHNICAL_SKILLS",     "")

    # shapes[6]: Name – Role header
    if len(shapes) > 6 and shapes[6].has_text_frame:
        header = f"{name} – {role}" if name and role else (name or role or "")
        tf = shapes[6].text_frame
        if tf.paragraphs:
            _set_para_text(tf.paragraphs[0], header)

    # shapes[1]: Professional summary
    if len(shapes) > 1 and shapes[1].has_text_frame:
        tf = shapes[1].text_frame
        if tf.paragraphs:
            _set_para_text(tf.paragraphs[0], summary or "")
            for p in tf.paragraphs[1:]:
                _clear_para(p)

    # shapes[2]: Education
    if len(shapes) > 2 and shapes[2].has_text_frame:
        tf = shapes[2].text_frame
        if tf.paragraphs:
            _set_para_text(
                tf.paragraphs[0],
                edu or "",
                color=BLACK if edu else RED,
            )
            for p in tf.paragraphs[1:]:
                _clear_para(p)

    # shapes[4]: Technical skills
    if len(shapes) > 4 and shapes[4].has_text_frame:
        tf = shapes[4].text_frame
        if tf.paragraphs:
            _set_para_text(tf.paragraphs[0], skills or "")
            for p in tf.paragraphs[1:]:
                _clear_para(p)


def _populate_project_slide(slide, proj_num: int, d: dict):
    shapes = slide.shapes

    proj_name  = d.get(f"PROJECT{proj_num}_NAME",        "")
    duration   = d.get(f"PROJECT{proj_num}_DURATION",    "")
    role       = d.get(f"PROJECT{proj_num}_ROLE",        "")
    desc       = d.get(f"PROJECT{proj_num}_DESCRIPTION", "")
    resp_lines = d.get(f"PROJECT{proj_num}_RESPONSIBILITIES", [])

    # Normalise resp_lines — _build_template_dict always sets a list,
    # but handle string just in case
    if isinstance(resp_lines, str):
        resp_lines = [
            re.sub(r"^[•\-\*]\s*", "", l).strip()
            for l in resp_lines.split("\n") if l.strip()
        ]
    elif not isinstance(resp_lines, list):
        resp_lines = []

    # shapes[1]: Project Name / Duration / Role
    if len(shapes) > 1 and shapes[1].has_text_frame:
        tf    = shapes[1].text_frame
        paras = tf.paragraphs

        def _lp(para, label, value):
            _set_para_text(
                para,
                f"{label}: {value}" if value else f"{label}: ",
                color=BLACK if value else RED,
            )

        if len(paras) > 0: _lp(paras[0], "Project Name", proj_name)
        if len(paras) > 1: _lp(paras[1], "Duration",     duration)
        if len(paras) > 2: _lp(paras[2], "Role",         role)

    # shapes[2]: Description + Responsibilities
    if len(shapes) > 2 and shapes[2].has_text_frame:
        tf    = shapes[2].text_frame
        paras = tf.paragraphs

        # para[0] = description
        if paras:
            _set_para_text(paras[0], desc or "", color=BLACK if desc else RED)

        # Find "Responsibilities:" label
        ri = _find_resp_label_idx(tf)
        if ri == -1:
            ri = min(4, len(paras) - 1)

        # Clear intermediate blank paragraphs
        for i in range(1, ri):
            _clear_para(paras[i])

        # Fill bullets after the label
        bs = ri + 1
        for i, line in enumerate(resp_lines):
            idx = bs + i
            if idx < len(paras):
                _set_para_text(paras[idx], line, color=BLACK)
            else:
                _append_para(tf, line)

        # Clear leftover old bullet slots
        for i in range(bs + len(resp_lines), len(paras)):
            _clear_para(paras[i])

# Public entry point

def generate_candidate_ppt(candidate_data: dict) -> bytes | None:
    """
    Generate a filled PPT from the NexTurn template.
    Accepts candidate_data in ANY key format (A/B/C/D or mixed).
    Returns raw bytes or None on failure.
    """
    try:
        if not os.path.exists(TEMPLATE_PATH):
            import streamlit as st
            st.error(
                f"PPT template not found at: {TEMPLATE_PATH}\n"
                "Place sample_ppt_template.pptx in the templates/ folder."
            )
            return None

        d   = _build_template_dict(candidate_data)
        prs = Presentation(TEMPLATE_PATH)
        n   = len(prs.slides)

        if n >= 1:
            _populate_slide1(prs.slides[0], d)

        for i in range(1, min(5, n)):
            _populate_project_slide(prs.slides[i], i, d)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    except Exception as e:
        import traceback
        import streamlit as st
        st.error(f"PPT Generation Error: {e}\n{traceback.format_exc()}")
        return None