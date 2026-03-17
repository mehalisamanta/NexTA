"""
utils/ppt_generator.py
Generates candidate PPT from the NexTurn template.

"""

import io
import os

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE

from utils.template_mapper import map_to_template_format
from utils.debug_log import debug_log


TEMPLATE_PATH = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "templates", "sample_ppt_template.pptx"
)

RED   = RGBColor(0xFF, 0x00, 0x00)
BLACK = RGBColor(0x00, 0x00, 0x00)


# Helpers 

def _enable_autofit(shape):
    """Shrink text to fit the text box automatically."""
    if shape.has_text_frame:
        tf = shape.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.word_wrap = True
        for para in tf.paragraphs:
            for run in para.runs:
                if run.font.size is None:
                    run.font.size = Pt(12)


def _apply_style(run, missing: bool):
    """Red bold for missing info, black for valid values."""
    if missing:
        run.font.color.rgb = RED
        run.font.bold = True
    else:
        run.font.color.rgb = BLACK


def _replace_placeholder(para, placeholder: str, value: str):
    """
    Replace {PLACEHOLDER} in a paragraph with value.
    Styles the run red if value is missing, black otherwise.
    Handles cases where placeholder is split across multiple runs.
    """
    # First try fast path — placeholder in a single run
    full_text = para.text
    if placeholder not in full_text:
        return

    content = str(value).strip() if value and str(value).strip() else "Information Missing"
    is_missing = (content == "Information Missing")

    # Rebuild paragraph text by replacing in the combined text,
    # then write it back into the first run and clear the rest
    new_text = full_text.replace(placeholder, content)

    if para.runs:
        para.runs[0].text = new_text
        _apply_style(para.runs[0], is_missing)
        for run in para.runs[1:]:
            run.text = ""
    else:
        # No runs — set text via XML (rare edge case)
        from pptx.oxml.ns import qn
        from lxml import etree
        r = para._p.get_or_add_r()
        r.text = new_text


# Slide 1: Profile 

def _populate_slide1(slide, d):
    placeholders = {
        "{FULL_NAME}":            d.get("FULL_NAME", ""),
        "{CURRENT_ROLE}":         d.get("CURRENT_ROLE", ""),
        "{PROFESSIONAL_SUMMARY}": d.get("PROFESSIONAL_SUMMARY", ""),
        "{EDUCATION_DETAILS}":    d.get("EDUCATION_DETAILS", ""),
        "{TECHNICAL_SKILLS}":     d.get("TECHNICAL_SKILLS", ""),
    }

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        _enable_autofit(shape)
        for para in shape.text_frame.paragraphs:
            for placeholder, value in placeholders.items():
                if placeholder in para.text:
                    _replace_placeholder(para, placeholder, value)


# ── Slides 2-5: Project slides ────────────────────────────────────────────────

def _populate_project_slide(slide, proj_num: int, d):
    placeholders = {
        f"{{PROJECT{proj_num}_NAME}}":             d.get(f"PROJECT{proj_num}_NAME", ""),
        f"{{PROJECT{proj_num}_DURATION}}":         d.get(f"PROJECT{proj_num}_DURATION", ""),
        f"{{PROJECT{proj_num}_ROLE}}":             d.get(f"PROJECT{proj_num}_ROLE", ""),
        f"{{PROJECT{proj_num}_DESCRIPTION}}":      d.get(f"PROJECT{proj_num}_DESCRIPTION", ""),
        f"{{PROJECT{proj_num}_RESPONSIBILITIES}}": d.get(f"PROJECT{proj_num}_RESPONSIBILITIES", ""),
    }

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        _enable_autofit(shape)
        tf = shape.text_frame

        for para in tf.paragraphs:
            # Handle RESPONSIBILITIES separately — multi-line bullet expansion
            resp_key = f"{{PROJECT{proj_num}_RESPONSIBILITIES}}"
            if resp_key in para.text:
                resp_value = placeholders[resp_key]
                if resp_value:
                    lines = [l.strip() for l in resp_value.split("\n") if l.strip()]
                    if lines:
                        # Write first line into this paragraph
                        content = lines[0]
                        if para.runs:
                            para.runs[0].text = content
                            _apply_style(para.runs[0], False)
                            for r in para.runs[1:]:
                                r.text = ""
                        # Add subsequent lines as new paragraphs
                        for line in lines[1:]:
                            new_p = tf.add_paragraph()
                            new_p.text = line
                            new_p.level = 0
                else:
                    _replace_placeholder(para, resp_key, "")
                continue

            # All other placeholders — straight replacement
            for placeholder, value in placeholders.items():
                if placeholder in para.text:
                    _replace_placeholder(para, placeholder, value)


# Main 

def generate_candidate_ppt(candidate_data: dict) -> bytes | None:
    """Generate candidate PPT from template. Returns bytes or None on failure."""
    try:
        if not os.path.exists(TEMPLATE_PATH):
            import streamlit as st
            st.error(
                f"PPT template not found. "
                f"Place `sample_ppt_template.pptx` in the `templates/` folder. "
                f"Expected: {TEMPLATE_PATH}"
            )
            return None

        prs = Presentation(TEMPLATE_PATH)
        d   = map_to_template_format(candidate_data)

        debug_log(
            location="utils/ppt_generator.py:generate_candidate_ppt",
            message="PPT generation input mapped",
            hypothesis_id="H2",
            data={
                "input_keys_sample": list((candidate_data or {}).keys())[:25],
                "mapped_full_name": (d.get("FULL_NAME") or "")[:80],
                "mapped_role": (d.get("CURRENT_ROLE") or "")[:80],
                "mapped_proj1": (d.get("PROJECT1_NAME") or "")[:80],
                "mapped_proj1_resp_len": len(d.get("PROJECT1_RESPONSIBILITIES") or ""),
            },
        )

        if len(prs.slides) >= 1:
            _populate_slide1(prs.slides[0], d)

        for i in range(1, min(5, len(prs.slides))):
            _populate_project_slide(prs.slides[i], i, d)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    except Exception as e:
        import streamlit as st
        st.error(f"PPT Generation Error: {e}")
        return None